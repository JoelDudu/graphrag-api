"""
Processador de arquivos Office e outros formatos
Suporta: PDF, DOCX, DOC, XLSX, XLS, PPTX, PPT, TXT, CSV
"""

import os
from pathlib import Path
from typing import List, Dict, Any
import logging

# Importações para diferentes tipos de arquivo
import pypdf
from docx import Document
import openpyxl
from pptx import Presentation

logger = logging.getLogger(__name__)


class FileProcessor:
    """Processa diferentes tipos de arquivos e extrai texto"""
    
    SUPPORTED_EXTENSIONS = {
        '.pdf': 'PDF Document',
        '.docx': 'Word Document',
        '.doc': 'Word Document (Legacy)',
        '.xlsx': 'Excel Spreadsheet',
        '.xls': 'Excel Spreadsheet (Legacy)',
        '.pptx': 'PowerPoint Presentation',
        '.ppt': 'PowerPoint Presentation (Legacy)',
        '.txt': 'Text File',
        '.csv': 'CSV File',
    }
    
    @classmethod
    def is_supported(cls, file_path: str) -> bool:
        """Verifica se o arquivo é suportado"""
        ext = Path(file_path).suffix.lower()
        return ext in cls.SUPPORTED_EXTENSIONS
    
    @classmethod
    def get_file_type(cls, file_path: str) -> str:
        """Retorna o tipo do arquivo"""
        ext = Path(file_path).suffix.lower()
        return cls.SUPPORTED_EXTENSIONS.get(ext, 'Unknown')
    
    @classmethod
    def extract_text(cls, file_path: str) -> str:
        """Extrai texto do arquivo baseado na extensão"""
        ext = Path(file_path).suffix.lower()
        
        if not cls.is_supported(file_path):
            raise ValueError(f"Tipo de arquivo não suportado: {ext}")
        
        try:
            if ext == '.pdf':
                return cls._extract_from_pdf(file_path)
            elif ext == '.docx':
                return cls._extract_from_docx(file_path)
            elif ext in ['.xlsx', '.xls']:
                return cls._extract_from_excel(file_path)
            elif ext == '.pptx':
                return cls._extract_from_pptx(file_path)
            elif ext in ['.txt', '.csv']:
                return cls._extract_from_text(file_path)
            elif ext == '.doc':
                # Para .doc antigo, tenta usar unstructured
                return cls._extract_from_doc_legacy(file_path)
            elif ext == '.ppt':
                # Para .ppt antigo, tenta usar unstructured
                return cls._extract_from_ppt_legacy(file_path)
            else:
                raise ValueError(f"Processador não implementado para: {ext}")
                
        except Exception as e:
            logger.error(f"Erro ao processar arquivo {file_path}: {str(e)}")
            raise
    
    @staticmethod
    def _extract_from_pdf(file_path: str) -> str:
        """Extrai texto de PDF"""
        text_parts = []
        with open(file_path, 'rb') as file:
            pdf_reader = pypdf.PdfReader(file)
            for page_num, page in enumerate(pdf_reader.pages):
                text = page.extract_text()
                if text.strip():
                    text_parts.append(f"--- Página {page_num + 1} ---\n{text}")
        
        return "\n\n".join(text_parts)
    
    @staticmethod
    def _extract_from_docx(file_path: str) -> str:
        """Extrai texto de DOCX"""
        doc = Document(file_path)
        text_parts = []
        
        # Extrair parágrafos
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        
        # Extrair tabelas
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join([cell.text.strip() for cell in row.cells])
                if row_text.strip():
                    text_parts.append(row_text)
        
        return "\n\n".join(text_parts)
    
    @staticmethod
    def _extract_from_excel(file_path: str) -> str:
        """Extrai texto de Excel"""
        workbook = openpyxl.load_workbook(file_path, data_only=True)
        text_parts = []
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text_parts.append(f"=== Planilha: {sheet_name} ===")
            
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                if row_text.strip():
                    text_parts.append(row_text)
            
            text_parts.append("")  # Linha em branco entre planilhas
        
        return "\n".join(text_parts)
    
    @staticmethod
    def _extract_from_pptx(file_path: str) -> str:
        """Extrai texto de PowerPoint"""
        prs = Presentation(file_path)
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            text_parts.append(f"=== Slide {slide_num} ===")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)
            
            text_parts.append("")  # Linha em branco entre slides
        
        return "\n\n".join(text_parts)
    
    @staticmethod
    def _extract_from_text(file_path: str) -> str:
        """Extrai texto de arquivo TXT ou CSV"""
        encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as file:
                    return file.read()
            except UnicodeDecodeError:
                continue
        
        # Se nenhum encoding funcionar, tenta com errors='ignore'
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
            return file.read()
    
    @staticmethod
    def _extract_from_doc_legacy(file_path: str) -> str:
        """Extrai texto de DOC antigo usando unstructured"""
        try:
            from unstructured.partition.doc import partition_doc
            elements = partition_doc(filename=file_path)
            return "\n\n".join([str(el) for el in elements])
        except Exception as e:
            logger.warning(f"Não foi possível processar .doc com unstructured: {e}")
            raise ValueError("Formato .doc não suportado. Por favor, converta para .docx")
    
    @staticmethod
    def _extract_from_ppt_legacy(file_path: str) -> str:
        """Extrai texto de PPT antigo usando unstructured"""
        try:
            from unstructured.partition.ppt import partition_ppt
            elements = partition_ppt(filename=file_path)
            return "\n\n".join([str(el) for el in elements])
        except Exception as e:
            logger.warning(f"Não foi possível processar .ppt com unstructured: {e}")
            raise ValueError("Formato .ppt não suportado. Por favor, converta para .pptx")
    
    @classmethod
    def get_supported_formats(cls) -> Dict[str, str]:
        """Retorna lista de formatos suportados"""
        return cls.SUPPORTED_EXTENSIONS.copy()
